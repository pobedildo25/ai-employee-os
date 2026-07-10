from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from app.document_intelligence.ast.models import ASTNode, ASTNodeType
from app.document_renderer.builders.style_applier import StyleApplier
from app.document_renderer.exceptions import RenderValidationError
from app.document_renderer.interfaces.renderer import DocumentRenderer
from app.document_renderer.models import OutputFormat, RenderRequest, RenderResult, RenderStatus


class PptxRenderer(DocumentRenderer):
    MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def __init__(self, style_applier: StyleApplier | None = None) -> None:
        self._style_applier = style_applier or StyleApplier()

    def validate(self, request: RenderRequest) -> None:
        if request.output_format != OutputFormat.PPTX:
            raise RenderValidationError("PptxRenderer only supports PPTX output")
        if request.document_structure.root.node_type != ASTNodeType.DOCUMENT:
            raise RenderValidationError("Document AST root must be a document node")

    def render(self, request: RenderRequest) -> RenderResult:
        self.validate(request)
        presentation = Presentation()
        self._style_applier.apply_pptx_slide_dimensions(presentation, request.brand_profile)

        sections = [
            child
            for child in request.document_structure.root.children
            if child.node_type == ASTNodeType.SECTION
        ]
        if not sections:
            sections = [request.document_structure.root]

        for section in sections:
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            self._render_section(slide, section, request)

        buffer = BytesIO()
        presentation.save(buffer)
        file_bytes = buffer.getvalue()

        return RenderResult(
            mime_type=self.MIME_TYPE,
            status=RenderStatus.COMPLETED,
            file_bytes=file_bytes,
            metadata={
                "format": OutputFormat.PPTX.value,
                "slide_count": len(presentation.slides),
            },
        )

    def _render_section(self, slide: Any, section: ASTNode, request: RenderRequest) -> None:
        top = Inches(1)
        for child in section.children:
            if child.node_type == ASTNodeType.HEADING:
                textbox = slide.shapes.add_textbox(Inches(1), top, Inches(8), Inches(1))
                textbox.text_frame.text = child.content or ""
                if textbox.text_frame.paragraphs and textbox.text_frame.paragraphs[0].runs:
                    self._style_applier.apply_pptx_text_style(
                        textbox.text_frame.paragraphs[0].runs[0],
                        brand_profile=request.brand_profile,
                        heading=True,
                    )
                top = Inches(2)
            elif child.node_type == ASTNodeType.PARAGRAPH:
                textbox = slide.shapes.add_textbox(Inches(1), top, Inches(8), Inches(1.5))
                textbox.text_frame.text = child.content or ""
                if textbox.text_frame.paragraphs and textbox.text_frame.paragraphs[0].runs:
                    self._style_applier.apply_pptx_text_style(
                        textbox.text_frame.paragraphs[0].runs[0],
                        brand_profile=request.brand_profile,
                    )
                top = Inches(3.5)
            elif child.node_type == ASTNodeType.IMAGE:
                image_data = child.attributes.get("image_bytes")
                if image_data:
                    slide.shapes.add_picture(BytesIO(image_data), Inches(1), top, width=Inches(4))
