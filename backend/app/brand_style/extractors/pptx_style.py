import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation

from app.brand_style.interfaces.extractor import StyleExtractor
from app.document_intelligence.models import DocumentRepresentation


class PptxStyleExtractor(StyleExtractor):
    def extract(
        self,
        document_representation: DocumentRepresentation,
        *,
        file_bytes: bytes | None = None,
        filename: str | None = None,
    ) -> dict[str, Any]:
        if not file_bytes:
            return self._extract_from_representation(document_representation)

        suffix = Path(filename or "document.pptx").suffix or ".pptx"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name

        try:
            presentation = Presentation(tmp_path)
        finally:
            Path(tmp_path).unlink(missing_ok=True)

        fonts: set[str] = set()
        font_sizes: set[float] = set()
        text_colors: set[str] = set()
        element_positions: list[dict[str, int]] = []
        slide_backgrounds: list[str] = []

        for slide in presentation.slides:
            background = slide.background
            if background.fill.type is not None:
                slide_backgrounds.append(str(background.fill.type))

            for shape in slide.shapes:
                element_positions.append(
                    {
                        "left": int(shape.left),
                        "top": int(shape.top),
                        "width": int(shape.width),
                        "height": int(shape.height),
                    }
                )
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
                        if run.font.size is not None:
                            font_sizes.add(float(run.font.size.pt))
                        if run.font.color.type is not None:
                            text_colors.add(str(run.font.color.type))

        theme_colors = self._extract_theme_colors(presentation)

        return {
            "format": "pptx",
            "slide_width": int(presentation.slide_width),
            "slide_height": int(presentation.slide_height),
            "slide_count": len(presentation.slides),
            "fonts": sorted(fonts),
            "font_sizes": sorted(font_sizes),
            "text_colors": sorted(text_colors),
            "theme_colors": theme_colors,
            "element_positions": element_positions[:20],
            "slide_backgrounds": slide_backgrounds,
        }

    def _extract_theme_colors(self, presentation: Presentation) -> dict[str, str]:
        theme_colors: dict[str, str] = {}
        try:
            theme_part = presentation.slide_master.theme_part
            theme = theme_part.theme
            color_scheme = theme.element.clrScheme
            for color_name in (
                "dk1",
                "lt1",
                "dk2",
                "lt2",
                "accent1",
                "accent2",
                "accent3",
                "accent4",
                "accent5",
                "accent6",
            ):
                color_element = getattr(color_scheme, color_name, None)
                if color_element is None:
                    continue
                srgb = getattr(getattr(color_element, "srgbClr", None), "val", None)
                if srgb is not None:
                    theme_colors[color_name] = f"#{srgb}"
        except Exception:
            theme_colors = {}
        return theme_colors

    def _extract_from_representation(self, representation: DocumentRepresentation) -> dict[str, Any]:
        return {
            "format": "pptx",
            "slide_count": representation.structure.get("pages", 0),
            "section_count": representation.structure.get("section_count", 0),
            "fonts": representation.metadata.get("fonts", []),
        }
