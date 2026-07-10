from pptx import Presentation

from app.file_processing.interfaces.extractor import Extractor
from app.file_processing.models import ExtractedContent


class PptxExtractor(Extractor):
    def extract(self, file_path: str) -> ExtractedContent:
        presentation = Presentation(file_path)
        slides: list[dict[str, object]] = []
        slide_texts: list[str] = []

        for index, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text.strip())
            slide_text = "\n".join(texts).strip()
            if slide_text:
                slide_texts.append(slide_text)
            slides.append({"index": index, "text": slide_text, "shape_count": len(slide.shapes)})

        full_text = "\n\n".join(slide_texts).strip()

        return ExtractedContent(
            text=full_text or None,
            metadata={"slide_count": len(presentation.slides), "format": "pptx"},
            pages=len(presentation.slides),
            structure={"slides": slides},
        )
