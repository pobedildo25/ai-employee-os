"""Generate test files for file processing tests."""

from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter


def create_txt_file(path: Path, content: str = "Hello from TXT file.\nLine two.") -> Path:
    path.write_text(content, encoding="utf-8")
    return path


def create_pdf_file(path: Path, text: str = "Sample PDF content") -> Path:
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with path.open("wb") as f:
        writer.write(f)
    return path


def create_docx_file(path: Path, paragraphs: list[str] | None = None) -> Path:
    document = Document()
    for paragraph in paragraphs or ["First paragraph", "Second paragraph"]:
        document.add_paragraph(paragraph)
    document.save(path)
    return path


def create_pptx_file(path: Path, slide_texts: list[str] | None = None) -> Path:
    presentation = Presentation()
    for text in slide_texts or ["Slide 1 title", "Slide 2 content"]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        textbox.text_frame.text = text
    presentation.save(path)
    return path


def create_xlsx_file(path: Path) -> Path:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sheet1"
    sheet.append(["Name", "Value"])
    sheet.append(["Item A", "100"])
    sheet.append(["Item B", "200"])
    workbook.save(path)
    return path


def create_png_file(path: Path) -> Path:
    # Minimal valid 1x1 PNG
    png_bytes = (
        b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00"
        b"\x01\x01\x01\x00\x18\xdd\x8d\xb4\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    path.write_bytes(png_bytes)
    return path


def create_large_pdf_file(path: Path, pages: int = 100) -> Path:
    writer = PdfWriter()
    for _ in range(pages):
        writer.add_blank_page(width=200, height=200)
    with path.open("wb") as f:
        writer.write(f)
    return path


def create_docx_with_table(path: Path) -> Path:
    document = Document()
    document.add_heading("Competitor Overview", level=1)
    table = document.add_table(rows=1, cols=3)
    headers = table.rows[0].cells
    headers[0].text = "Company"
    headers[1].text = "Strength"
    headers[2].text = "Weakness"
    for row in [("Alpha", "Brand", "Price"), ("Beta", "Speed", "Scale")]:
        cells = table.add_row().cells
        for index, value in enumerate(row):
            cells[index].text = value
    document.save(path)
    return path


def create_pptx_with_image(path: Path, image_path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1)).text_frame.text = "Branded deck"
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), width=Inches(2))
    presentation.save(path)
    return path


def read_file_bytes(path: Path) -> bytes:
    return path.read_bytes()
