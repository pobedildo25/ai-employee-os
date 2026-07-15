from io import BytesIO
from uuid import UUID, uuid4

import pytest
from docx import Document
from pptx import Presentation

from app.brand_style.models import BrandProfile
from app.document_intelligence.ast.models import ASTNode, ASTNodeType, DocumentAST
from app.document_renderer.builders.style_applier import StyleApplier
from app.document_renderer.exceptions import UnsupportedFormatError
from app.document_renderer.models import OutputFormat, RenderRequest, RenderStatus
from app.document_renderer.renderer import DocumentRendererService, RenderArtifactService
from app.skills.builtin.document_render_skill import DocumentRenderSkill
from app.skills.registry import create_capability_registry


def _sample_ast(*, for_pptx: bool = False) -> DocumentAST:
    if for_pptx:
        section = ASTNode(
            node_type=ASTNodeType.SECTION,
            content="Slide 1",
            children=[
                ASTNode(node_type=ASTNodeType.HEADING, content="Presentation Title"),
                ASTNode(node_type=ASTNodeType.PARAGRAPH, content="Slide body text"),
            ],
        )
    else:
        section = ASTNode(
            node_type=ASTNodeType.SECTION,
            content="Body",
            children=[
                ASTNode(node_type=ASTNodeType.HEADING, content="Document Title"),
                ASTNode(node_type=ASTNodeType.PARAGRAPH, content="First paragraph content"),
                ASTNode(
                    node_type=ASTNodeType.TABLE,
                    content="Data Table",
                    attributes={"rows": [["Name", "Value"], ["Item A", "100"]]},
                ),
            ],
        )

    root = ASTNode(
        node_type=ASTNodeType.DOCUMENT,
        content="Generated Document",
        children=[section],
    )
    return DocumentAST(root=root, node_count=5)


@pytest.fixture
def brand_profile() -> BrandProfile:
    return BrandProfile(
        client_id=uuid4(),
        name="Test Brand",
        typography={
            "body_font": "Calibri",
            "heading_font": "Arial",
            "font_sizes": [11.0, 18.0],
        },
        colors={"primary": "#123456"},
        layout_rules={
            "header": "top",
            "footer": True,
            "page_margins": {"top": 72.0, "bottom": 72.0, "left": 72.0, "right": 72.0},
            "slide_dimensions": {"width": 9144000, "height": 6858000},
        },
    )


@pytest.fixture
def renderer_service() -> DocumentRendererService:
    return DocumentRendererService()


def test_render_request_example(brand_profile: BrandProfile) -> None:
    request = RenderRequest(
        document_structure=_sample_ast(),
        brand_profile=brand_profile,
        output_format=OutputFormat.DOCX,
        metadata={"title": "Sample Output"},
        client_id=uuid4(),
        project_id=uuid4(),
        name="output.docx",
    )

    assert request.output_format == OutputFormat.DOCX
    assert request.brand_profile.typography["body_font"] == "Calibri"


def test_docx_generation(renderer_service: DocumentRendererService, brand_profile: BrandProfile) -> None:
    request = RenderRequest(
        document_structure=_sample_ast(),
        brand_profile=brand_profile,
        output_format=OutputFormat.DOCX,
        metadata={"title": "DOCX Test"},
    )

    result = renderer_service.render(request)

    assert result.status == RenderStatus.COMPLETED
    assert result.file_bytes is not None
    assert result.mime_type.endswith("wordprocessingml.document")
    assert result.metadata.get("template") == "herald"

    document = Document(BytesIO(result.file_bytes))
    texts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    assert "DOCX Test" in texts  # centered HERALD title
    assert "Document Title" in texts
    assert "First paragraph content" in texts
    assert len(document.tables) == 1
    # Header chrome table (mark + wordmark) + body data table.
    header = document.sections[0].header
    assert header.tables
    assert any(rel.reltype.endswith("/image") for rel in header.part.rels.values())


def test_pptx_generation(renderer_service: DocumentRendererService, brand_profile: BrandProfile) -> None:
    request = RenderRequest(
        document_structure=_sample_ast(for_pptx=True),
        brand_profile=brand_profile,
        output_format=OutputFormat.PPTX,
    )

    result = renderer_service.render(request)

    assert result.status == RenderStatus.COMPLETED
    assert result.file_bytes is not None

    presentation = Presentation(BytesIO(result.file_bytes))
    assert len(presentation.slides) == 1
    slide_texts: list[str] = []
    for shape in presentation.slides[0].shapes:
        if hasattr(shape, "text") and shape.text:
            slide_texts.append(shape.text)
    assert any("Presentation Title" in text for text in slide_texts)
    assert any("Slide body text" in text for text in slide_texts)


def test_style_applier_uses_brand_profile(brand_profile: BrandProfile) -> None:
    applier = StyleApplier()
    assert applier.get_body_font(brand_profile) == "Calibri"
    assert applier.get_heading_font(brand_profile) == "Arial"
    assert applier.get_primary_color(brand_profile) == "#123456"


def test_pdf_renderer_not_implemented(renderer_service: DocumentRendererService, brand_profile: BrandProfile) -> None:
    request = RenderRequest(
        document_structure=_sample_ast(),
        brand_profile=brand_profile,
        output_format=OutputFormat.PDF,
    )

    with pytest.raises(UnsupportedFormatError):
        renderer_service.render(request)


@pytest.mark.asyncio
async def test_artifact_creation(brand_profile: BrandProfile, artifact_service) -> None:
    client_id, project_id = uuid4(), uuid4()
    render_service = RenderArtifactService(artifact_service=artifact_service)

    request = RenderRequest(
        document_structure=_sample_ast(),
        brand_profile=brand_profile,
        output_format=OutputFormat.DOCX,
        client_id=client_id,
        project_id=project_id,
        name="generated.docx",
        brand_profile_id=brand_profile.id,
        source_artifact_id=uuid4(),
        metadata={"title": "Artifact Test"},
    )

    result = await render_service.render_and_store(request)

    assert result.artifact_id is not None
    assert result.file_path is not None
    assert result.metadata.get("generated_by") == "document_renderer"
    assert result.metadata.get("brand_profile_id") == str(brand_profile.id)

    artifact = await artifact_service.get_by_id(result.artifact_id)
    assert artifact is not None
    assert artifact.artifact_type == "generated_document"
    assert artifact.metadata is not None
    assert artifact.metadata.get("output_format") == "docx"


@pytest.mark.asyncio
async def test_document_render_skill(brand_profile: BrandProfile) -> None:
    skill = DocumentRenderSkill()
    result = await skill.execute(
        {
            "document_ast": _sample_ast().model_dump(mode="json"),
            "brand_profile": brand_profile.model_dump(mode="json"),
            "output_format": "docx",
        }
    )

    assert result["status"] == "completed"
    render_result = result["render_result"]
    assert render_result["status"] == RenderStatus.COMPLETED.value
    assert render_result["file_size"] > 0


@pytest.mark.asyncio
async def test_document_render_skill_with_artifact(
    brand_profile: BrandProfile,
    artifact_service,
) -> None:
    client_id, project_id = uuid4(), uuid4()
    skill = DocumentRenderSkill(artifact_service=RenderArtifactService(artifact_service=artifact_service))

    result = await skill.execute(
        {
            "document_ast": _sample_ast(for_pptx=True).model_dump(mode="json"),
            "brand_profile": brand_profile.model_dump(mode="json"),
            "output_format": "pptx",
            "client_id": str(client_id),
            "project_id": str(project_id),
            "name": "generated.pptx",
            "store_artifact": True,
        }
    )

    assert result["status"] == "completed"
    render_result = result["render_result"]
    assert render_result["artifact_id"] is not None

    artifact = await artifact_service.get_by_id(UUID(str(render_result["artifact_id"])))
    assert artifact is not None
    assert artifact.storage_path is not None
    file_bytes = await artifact_service._storage.download(artifact.storage_path)
    presentation = Presentation(BytesIO(file_bytes))
    assert len(presentation.slides) >= 1


def test_skill_registry_includes_document_rendering() -> None:
    registry = create_capability_registry()
    names = {capability.name for capability in registry.list_available()}
    assert "document_rendering" in names
    skill = registry.get_skill_for_capability("document_rendering")
    assert skill is not None
    assert skill.name() == "document_render_skill"
